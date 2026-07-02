#!/usr/bin/env python3
"""AnantaSutra pitch deck (.pptx) — Dedicated Domain Experts model.
Built from real site content: dedicated-experts page, expertDomains, clients."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from PIL import Image

REPO = r"c:/Users/duneja8515/Desktop/bhavya/prsnl/Portfolio"
CLIENTS_DIR = os.path.join(REPO, "public", "images", "clients")
PORTRAIT = os.path.join(REPO, "public", "images", "portrait.jpg")

# ---- Brand palette ----
BG      = RGBColor(0x0A, 0x0A, 0x0F)
CARD    = RGBColor(0x16, 0x16, 0x20)
CARD2   = RGBColor(0x1B, 0x16, 0x12)
BORDER  = RGBColor(0x2A, 0x2A, 0x38)
SAFFRON = RGBColor(0xE8, 0xA3, 0x17)
SAFF_LT = RGBColor(0xF0, 0xC0, 0x40)
VIOLET  = RGBColor(0x6A, 0x3D, 0xE8)
VIO_LT  = RGBColor(0x9B, 0x7C, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT    = RGBColor(0xC9, 0xCA, 0xD3)
MUTED   = RGBColor(0x82, 0x82, 0x92)

DISP = "Segoe UI"
BODY = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    return s


def rect(s, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    r = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def oval_outline(s, cx, cy, d, color, weight=1.1):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    o.fill.background(); o.line.color.rgb = color; o.line.width = Pt(weight)
    o.shadow.inherit = False
    return o


def fit_image(s, path, bl, bt, bw, bh):
    iw, ih = Image.open(path).size
    scale = min(bw / iw, bh / ih)
    w_in, h_in = iw * scale, ih * scale
    s.shapes.add_picture(path, Inches(bl + (bw - w_in) / 2), Inches(bt + (bh - h_in) / 2),
                         Inches(w_in), Inches(h_in))


def square_crop(path, out):
    im = Image.open(path).convert("RGB")
    w, h = im.size; m = min(w, h)
    left = (w - m) // 2; top = (h - m) // 2
    im.crop((left, top, left + m, top + m)).save(out, quality=92)
    return out


def card(s, l, t, w, h, fill=CARD, border=BORDER):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = border; c.line.width = Pt(0.75)
    c.shadow.inherit = False
    try: c.adjustments[0] = 0.05
    except Exception: pass
    return c


def tb(s, l, t, w, h):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    return tf


def put(tf, text, size, color, bold=False, font=BODY, align=PP_ALIGN.LEFT,
        after=4, before=0, italic=False, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.space_before = Pt(before)
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font; f.color.rgb = color
    return p


def header(s, kicker, title, size=32):
    put(tb(s, 0.7, 0.5, 12.2, 0.35), kicker.upper(), 12, SAFFRON, bold=True, first=True)
    put(tb(s, 0.7, 0.82, 12.2, 1.0), title, size, WHITE, bold=True, font=DISP, first=True)
    rect(s, 0.73, 1.72, 0.85, 0.045, SAFFRON)


def footer(s, n):
    put(tb(s, 0.7, 7.04, 9, 0.3), "AnantaSutra   ·   anantasutra.com", 9, MUTED, first=True)
    put(tb(s, 11.9, 7.04, 0.9, 0.3), str(n), 9, MUTED, align=PP_ALIGN.RIGHT, first=True)


def cols(n, total=11.93, left0=0.7, gap=0.3):
    w = (total - (n - 1) * gap) / n
    return [left0 + i * (w + gap) for i in range(n)], w


def feature_card(s, l, t, w, h, title, lines, accent, title_color=WHITE):
    c = card(s, l, t, w, h)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.26); tf.margin_right = Inches(0.24)
    tf.margin_top = Inches(0.26); tf.margin_bottom = Inches(0.2)
    put(tf, title, 15.5, title_color, bold=True, font=DISP, after=7, first=True)
    for ln in lines:
        put(tf, ln, 11.5, TEXT, after=3)
    rect(s, l + 0.26, t + 0.2, 0.34, 0.035, accent)
    return c


def metric_card(s, l, t, w, h, value, name, sub, accent):
    c = card(s, l, t, w, h)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.26); tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.24); tf.margin_bottom = Inches(0.18)
    put(tf, value, 25, accent, bold=True, font=DISP, after=5, first=True)
    put(tf, name, 13, WHITE, bold=True, after=3)
    if sub:
        put(tf, sub, 10.5, TEXT, after=0)
    return c


def stat_box(s, l, t, w, h, value, label, accent):
    c = card(s, l, t, w, h)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    put(tf, value, 26, accent, bold=True, font=DISP, align=PP_ALIGN.CENTER, after=3, first=True)
    put(tf, label, 11, TEXT, align=PP_ALIGN.CENTER, after=0)


# ===================== SLIDE 1 — COVER =====================
s = slide()
for i, d in enumerate([5.0, 3.9, 2.8, 1.7]):
    oval_outline(s, 10.0, 3.75, d, SAFFRON if i % 2 == 0 else VIOLET)
put(tb(s, 8.7, 3.0, 2.6, 1.5), "ॐ", 54, RGBColor(0x3A, 0x2E, 0x12), bold=True, align=PP_ALIGN.CENTER, first=True)
put(tb(s, 0.9, 1.75, 7.5, 0.4), "अनन्तसूत्र   ·   ANANTASUTRA", 14, SAFFRON, bold=True, first=True)
put(tb(s, 0.86, 2.15, 8.0, 1.2), "AnantaSutra", 56, WHITE, bold=True, font=DISP, first=True)
rect(s, 0.92, 3.35, 1.1, 0.05, SAFFRON)
put(tb(s, 0.9, 3.55, 7.6, 0.6), "Not a project. A dedicated expert.", 22, SAFF_LT, bold=True, first=True)
put(tb(s, 0.9, 4.25, 7.1, 1.2),
    "Vetted domain professionals — any field — embedded inside your team, "
    "for as long as you need. You pay just the salary.",
    14.5, TEXT, first=True)
put(tb(s, 0.9, 6.25, 8.5, 0.4), "Bhavya Duneja — Co-founder    ·    Delhi, India", 12.5, WHITE, bold=True, first=True)
put(tb(s, 0.9, 6.62, 8.5, 0.4), "anantasutra.com    ·    contact@anantasutra.com", 12.5, SAFFRON, first=True)

# ===================== SLIDE 2 — THE CORE DIFFERENCE =====================
s = slide(); header(s, "What we do", "We don't sell deliverables. We place a person.")
tf = tb(s, 0.7, 2.2, 7.6, 3.4)
put(tf, "Most agencies hand you a project and disappear. AnantaSutra assigns a vetted expert "
        "who becomes part of your team — they attend your standups, learn your product, and own "
        "your KPIs, because they're not splitting time across ten clients.", 16.5, TEXT, after=14, first=True)
put(tf, "One domain or many. In-house or remote. An afternoon or a full year.", 16.5, TEXT, after=0)
q = card(s, 8.55, 2.25, 4.1, 3.3, fill=CARD2)
qtf = q.text_frame; qtf.vertical_anchor = MSO_ANCHOR.MIDDLE; qtf.word_wrap = True
qtf.margin_left = Inches(0.3); qtf.margin_right = Inches(0.28)
put(qtf, "“You run your business.", 19, WHITE, bold=True, font=DISP, after=2, first=True)
put(qtf, "Your expert runs your execution.”", 19, SAFF_LT, bold=True, font=DISP, after=0)
footer(s, 2)

# ===================== SLIDE 3 — PROBLEM =====================
s = slide(); header(s, "The problem", "Hiring is slow. Agencies are detached.")
xs, w = cols(3)
probs = [
    ("Hiring takes months", ["Job posts, interviews, notice periods.",
                             "By the time someone finally starts,",
                             "the opportunity has already passed."], SAFFRON),
    ("Agencies juggle ten clients", ["And mark talent up 30–50%. You pay",
                                     "₹1.5L for someone earning ₹1L —",
                                     "and you're just one account of many."], VIOLET),
    ("Freelancers aren't embedded", ["No ownership, no context, no real",
                                     "accountability. Work slips through",
                                     "the cracks between hand-offs."], SAFFRON),
]
for x, (t, lines, acc) in zip(xs, probs):
    feature_card(s, x, 2.25, w, 3.5, t, lines, acc)
footer(s, 3)

# ===================== SLIDE 4 — THE MODEL =====================
s = slide(); header(s, "Our model", "Dedicated domain experts, on demand.")
xs, w = cols(3)
sol = [
    ("Any Domain, One Partner", ["An engineer today, a lawyer next",
                                 "month — every field through a single",
                                 "relationship. No juggling vendors."], SAFFRON),
    ("Embedded, Not Outsourced", ["Your expert works inside your team,",
                                  "under your direction — in-house or",
                                  "remote. Real ownership, not tickets."], VIOLET),
    ("Flexible by Design", ["Scale up when you grow, wind down",
                            "when you don't. Add experts, pause",
                            "anytime — no lock-ins."], SAFFRON),
]
for x, (t, lines, acc) in zip(xs, sol):
    feature_card(s, x, 2.25, w, 3.5, t, lines, acc, title_color=SAFF_LT)
footer(s, 4)

# ===================== SLIDE 5 — ANY DOMAIN (12 domains) =====================
s = slide(); header(s, "Any domain", "Whatever expertise you need.")
domains = [
    ("Engineering & Software", "Full-Stack · Backend · DevOps · Mobile"),
    ("AI & Data", "ML · Data Science · Automation"),
    ("Design & Creative", "Product · Brand · UI/UX · Motion"),
    ("Marketing & Growth", "Performance · Growth · SEO · Social"),
    ("Video & Content", "Editors · Writers · Producers"),
    ("Legal & Compliance", "Corporate · Immigration · Contracts"),
    ("Healthcare", "Clinical Advisory · Health Ops · Telehealth"),
    ("Finance & Business", "Analysts · FP&A · Ops · Strategy"),
    ("Property & Real Estate", "Sales · Funnels · Listings · CRM"),
    ("E-commerce & Retail", "Store Ops · Catalogue · CRO · Retention"),
    ("Immigration & Mobility", "Case Mgmt · Visa · Documentation"),
    ("Academic & Research", "Researchers · Subject Experts · Writers"),
]
xs, w = cols(4, gap=0.26)
rows_top = [2.1, 3.52, 4.94]
ch = 1.3
for i, (name, roles) in enumerate(domains):
    row, col = i // 4, i % 4
    cl, ct = xs[col], rows_top[row]
    card(s, cl, ct, w, ch)
    acc = SAFFRON if i % 2 == 0 else VIOLET
    rect(s, cl + 0.2, ct + 0.22, 0.26, 0.03, acc)
    tf = tb(s, cl + 0.2, ct + 0.32, w - 0.36, ch - 0.4)
    put(tf, name, 12.5, WHITE, bold=True, after=3, first=True)
    put(tf, roles, 9.5, TEXT, after=0)
put(tb(s, 0.7, 6.4, 11.93, 0.4),
    "If it's a profession, we can embed an expert in it. The domain can be anything.",
    12.5, SAFF_LT, italic=True, align=PP_ALIGN.CENTER, first=True)
footer(s, 5)

# ===================== SLIDE 6 — HOW IT WORKS =====================
s = slide(); header(s, "How it works", "From need to onboarded in 7 days.")
xs, w = cols(4)
steps = [
    ("01", "Tell us your need", "A 30-minute call on your business, the role you need, and the outcomes you want."),
    ("02", "We match & vet", "We handpick 2–3 vetted pros from our network. You interview them and choose."),
    ("03", "They join your team", "Into your Slack, Notion and project tools — working as part of your team within 7 days."),
    ("04", "You stay focused", "Weekly reviews, monthly reports, and one accountable point of contact."),
]
for x, (num, title, desc) in zip(xs, steps):
    c = card(s, x, 2.25, w, 3.5)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.24); tf.margin_right = Inches(0.22); tf.margin_top = Inches(0.26)
    put(tf, num, 34, SAFF_LT, bold=True, font=DISP, after=4, first=True)
    put(tf, title, 15, WHITE, bold=True, font=DISP, after=8)
    put(tf, desc, 11.5, TEXT, after=2)
footer(s, 6)

# ===================== SLIDE 7 — WHY THIS MODEL =====================
s = slide(); header(s, "Why it works", "Built for founders who value time.")
why = [
    ("Dedicated, not shared", "Your expert works only on your business — never juggling five clients.", SAFFRON),
    ("Ramp-up in 7 days", "Skip the 3-month hiring cycle. Productive from week one.", VIOLET),
    ("Replace anytime", "Wrong fit? We replace them within 48 hours. No long contracts.", SAFFRON),
    ("Scale on demand", "Start with one expert. Add more as you grow. Pause when you don't.", VIOLET),
    ("Backed by AnantaSutra", "Every expert has our processes, tools and leadership behind them.", SAFFRON),
    ("Partnership mindset", "We grow when you grow — that's the whole reason this model exists.", VIOLET),
]
xs, w = cols(3)
rows_top = [2.15, 4.05]
ch = 1.75
for i, (title, desc, acc) in enumerate(why):
    row, col = i // 3, i % 3
    cl, ct = xs[col], rows_top[row]
    card(s, cl, ct, w, ch)
    rect(s, cl + 0.26, ct + 0.24, 0.3, 0.035, acc)
    tf = tb(s, cl + 0.26, ct + 0.36, w - 0.5, ch - 0.5)
    put(tf, title, 14, WHITE, bold=True, font=DISP, after=5, first=True)
    put(tf, desc, 11.5, TEXT, after=0)
footer(s, 7)

# ===================== SLIDE 8 — PRICING =====================
s = slide(); header(s, "Pricing", "You pay just the salary.")
xs, w = cols(3)
price = [
    ("What you pay", ["The expert's monthly salary —", "exactly what they earn.", "That's it."], SAFFRON, True),
    ("What you don't pay", ["No agency markup (save 30–50%).", "No recruitment fees. No overheads.", "No hidden charges."], VIOLET, False),
    ("How we earn", ["A small, transparent coordination", "fee — agreed upfront.", "You see every rupee."], SAFFRON, False),
]
for x, (title, lines, acc, hl) in zip(xs, price):
    c = card(s, x, 2.25, w, 3.05, fill=CARD2 if hl else CARD, border=acc if hl else BORDER)
    if hl:
        c.line.width = Pt(1.75)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.28); tf.margin_top = Inches(0.3)
    put(tf, title, 17, SAFF_LT if hl else WHITE, bold=True, font=DISP, after=10, first=True)
    for ln in lines:
        put(tf, ln, 13, TEXT, after=5)
put(tb(s, 0.7, 5.65, 11.93, 0.6),
    "Agencies charge ₹1.5L for someone earning ₹1L. We flipped it — you pay the salary, "
    "we earn a transparent fee, and you save lakhs every year.",
    12.5, MUTED, italic=True, align=PP_ALIGN.CENTER, first=True)
footer(s, 8)

# ===================== SLIDE 9 — TRACTION =====================
s = slide(); header(s, "Traction", "Numbers that speak.")
xs, w = cols(4)
strip = [("12+", "Expert Domains"), ("10+", "Brands Delivered"),
         ("4", "Continents Served"), ("7-Day", "Avg. Onboarding")]
for x, (v, l) in zip(xs, strip):
    stat_box(s, x, 2.25, w, 1.5, v, l, SAFFRON)
xs3, w3 = cols(3)
row2 = [("0%", "Agency markup"), ("48h", "Replacement guarantee"),
        ("100%", "Dedicated to your business")]
for x, (v, l) in zip(xs3, row2):
    stat_box(s, x, 4.0, w3, 1.5, v, l, VIOLET)
put(tb(s, 0.7, 5.7, 11.9, 0.5),
    "A vetted network spanning 12+ disciplines — embedded, accountable, and aligned to your outcomes.",
    13, TEXT, italic=True, align=PP_ALIGN.CENTER, first=True)
footer(s, 9)

# ===================== SLIDE 10 — CLIENT RESULTS =====================
s = slide(); header(s, "Client results", "Real brands. Real outcomes.")
xs, w = cols(4)
res = [
    ("₹25L → ₹40L", "Awish Clinic", "Monthly revenue in 2 months — embedded web, CRM & ads experts (Dermatology).", SAFFRON),
    ("₹23L", "BotWot", "Average deal size closed across India, UAE & Nigeria via an embedded marketing team.", VIOLET),
    ("+46%", "Zoom Wheels", "Revenue growth driven by a dedicated ads expert (Automotive, Delhi).", SAFFRON),
    ("+38%", "Royal Properties", "Traffic lift from embedded funnel & landing-page experts (Real Estate).", VIOLET),
]
for x, (v, n, sub, acc) in zip(xs, res):
    metric_card(s, x, 2.25, w, 3.5, v, n, sub, acc)
footer(s, 10)

# ===================== SLIDE 11 — TRUSTED BY (logo wall) =====================
s = slide(); header(s, "Trusted by", "Across industries and continents.")
roster = [
    ("Awish Clinic", "Dermatology · India", "awish-clinic.png"),
    ("Education Aspire", "EdTech · Faridabad", "education-aspire.png"),
    ("Giant Migrations", "Immigration · Qatar·UAE", "giant-migrations.png"),
    ("BotWot", "AI · India·UAE·Nigeria", "botwot.png"),
    ("Zoom Wheels", "Automotive · Delhi", "zoom-wheels.png"),
    ("Royal Properties", "Real Estate · Delhi", "royal-properties.png"),
    ("BlueMoon Marketing", "Advertising · Delhi", "bluemoon-marketing.png"),
    ("Smile With Kris", "Dental · UK", "smile-with-kris.png"),
    ("Walk Through My Lens", "Travel · London", "walk-through-my-lens.png"),
    ("Wisdom of Mind", "Vedic & Wellness · Haryana", "wisdom-of-mind.png"),
]
xs, w = cols(5, gap=0.28)
rows_top = [2.2, 4.35]
ch = 2.0
for i, (name, meta, logo) in enumerate(roster):
    row, col = i // 5, i % 5
    cl, ct = xs[col], rows_top[row]
    card(s, cl, ct, w, ch)
    chip_l, chip_t, chip_w, chip_h = cl + 0.16, ct + 0.16, w - 0.32, 0.86
    card(s, chip_l, chip_t, chip_w, chip_h, fill=WHITE, border=WHITE)
    lp = os.path.join(CLIENTS_DIR, logo)
    if os.path.exists(lp):
        fit_image(s, lp, chip_l + 0.13, chip_t + 0.13, chip_w - 0.26, chip_h - 0.26)
    acc = SAFFRON if i % 2 == 0 else VIOLET
    rect(s, cl + 0.16, ct + 1.18, 0.24, 0.03, acc)
    tf = tb(s, cl + 0.16, ct + 1.24, w - 0.28, 0.7)
    put(tf, name, 11.5, WHITE, bold=True, after=2, first=True)
    put(tf, meta, 9, TEXT, after=0)
footer(s, 11)

# ===================== SLIDE 12 — LEADERSHIP =====================
s = slide(); header(s, "Leadership", "Backed by a technologist-entrepreneur.")
_sq = square_crop(PORTRAIT, "_portrait_sq.jpg")
s.shapes.add_picture(_sq, Inches(0.7), Inches(2.3), Inches(2.9), Inches(2.9))
_frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.3), Inches(2.9), Inches(2.9))
_frame.fill.background(); _frame.line.color.rgb = SAFFRON; _frame.line.width = Pt(1.5); _frame.shadow.inherit = False
try: _frame.adjustments[0] = 0.04
except Exception: pass
tf = tb(s, 4.1, 2.35, 8.5, 3.2)
put(tf, "Bhavya Duneja", 26, WHITE, bold=True, font=DISP, after=2, first=True)
put(tf, "Co-founder, AnantaSutra", 14, SAFFRON, bold=True, after=12)
put(tf, "A software engineer and entrepreneur (.NET, React, AWS, Azure) bridging Delhi and "
        "Osaka, Japan. Bhavya leads a vetted network of professionals across 12+ domains — "
        "every expert backed by AnantaSutra's processes, tools and hands-on leadership.", 14, TEXT, after=12)
put(tf, "We grow when our clients grow — that's the whole reason this model exists.",
    13, SAFF_LT, italic=True, after=0)
footer(s, 12)

# ===================== SLIDE 13 — CLOSE / CONTACT =====================
s = slide()
for i, d in enumerate([5.4, 4.1, 2.8]):
    oval_outline(s, 6.66, 3.4, d, SAFFRON if i % 2 == 0 else VIOLET)
put(tb(s, 1, 1.7, 11.3, 0.4), "STOP JUGGLING. START GROWING.", 14, SAFFRON, bold=True, align=PP_ALIGN.CENTER, first=True)
put(tb(s, 1, 2.15, 11.3, 1.0), "Get your expert in 7 days.", 44, WHITE, bold=True, font=DISP, align=PP_ALIGN.CENTER, first=True)
put(tb(s, 1.5, 3.4, 10.3, 0.7),
    "Tell us the role you need — we'll introduce your dedicated expert within a week. "
    "Free consultation, zero commitment.",
    16, TEXT, align=PP_ALIGN.CENTER, first=True)
cta = card(s, 4.92, 4.4, 3.5, 0.7, fill=SAFFRON, border=SAFFRON)
ctf = cta.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
put(ctf, "contact@anantasutra.com", 15, BG, bold=True, align=PP_ALIGN.CENTER, first=True)
put(tb(s, 1, 5.55, 11.3, 0.4), "anantasutra.com    ·    Delhi, India", 13, TEXT, align=PP_ALIGN.CENTER, first=True)
put(tb(s, 1, 6.0, 11.3, 0.4), "Bhavya Duneja — Co-founder", 12, MUTED, align=PP_ALIGN.CENTER, first=True)

out = "AnantaSutra-Pitch-Deck.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
