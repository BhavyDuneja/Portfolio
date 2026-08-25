#!/usr/bin/env python3
"""AnantaSutra x Wonder Gifts — discovery-call deck (.pptx).
Same brand system as AnantaSutra-Pitch-Deck: saffron/violet on near-black."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Brand palette ----
BG      = RGBColor(0x0A, 0x0A, 0x0F)
CARD    = RGBColor(0x16, 0x16, 0x20)
CARD2   = RGBColor(0x1B, 0x16, 0x12)
BORDER  = RGBColor(0x2A, 0x2A, 0x38)
SAFFRON = RGBColor(0xE8, 0xA3, 0x17)
SAFF_LT = RGBColor(0xF0, 0xC0, 0x40)
VIOLET  = RGBColor(0x6A, 0x3D, 0xE8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT    = RGBColor(0xC9, 0xCA, 0xD3)
MUTED   = RGBColor(0x82, 0x82, 0x92)

DISP = "Segoe UI"
BODY = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    return s


def rect(s, l, t, w, h, color):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def oval_outline(s, cx, cy, d, color, weight=1.1):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    o.fill.background(); o.line.color.rgb = color; o.line.width = Pt(weight)
    o.shadow.inherit = False
    return o


def card(s, l, t, w, h, fill=CARD, border=BORDER):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = border; c.line.width = Pt(0.75)
    c.shadow.inherit = False
    try: c.adjustments[0] = 0.05
    except Exception: pass
    return c


def tb(s, l, t, w, h):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    return tf


def put(tf, text, size, color, bold=False, font=BODY, align=PP_ALIGN.LEFT,
        after=4, before=0, italic=False, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.space_before = Pt(before)
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font; f.color.rgb = color
    return p


def header(s, kicker, title, size=32):
    put(tb(s, 0.7, 0.5, 12.2, 0.35), kicker.upper(), 12, SAFFRON, bold=True, first=True)
    put(tb(s, 0.7, 0.82, 12.2, 1.0), title, size, WHITE, bold=True, font=DISP, first=True)
    rect(s, 0.73, 1.72, 0.85, 0.045, SAFFRON)


def footer(s, n):
    put(tb(s, 0.7, 7.04, 9, 0.3), "AnantaSutra   ·   anantasutra.com", 9, MUTED, first=True)
    put(tb(s, 11.9, 7.04, 0.9, 0.3), str(n), 9, MUTED, align=PP_ALIGN.RIGHT, first=True)


def cols(n, total=11.93, left0=0.7, gap=0.3):
    w = (total - (n - 1) * gap) / n
    return [left0 + i * (w + gap) for i in range(n)], w


def feature_card(s, l, t, w, h, title, lines, accent, title_color=WHITE, tsize=15.5, bsize=11.5):
    c = card(s, l, t, w, h)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.26); tf.margin_right = Inches(0.24)
    tf.margin_top = Inches(0.26); tf.margin_bottom = Inches(0.2)
    put(tf, title, tsize, title_color, bold=True, font=DISP, after=7, first=True)
    for ln in lines:
        put(tf, ln, bsize, TEXT, after=3)
    rect(s, l + 0.26, t + 0.2, 0.34, 0.035, accent)
    return c


def stat_box(s, l, t, w, h, value, label, accent):
    c = card(s, l, t, w, h)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    put(tf, value, 24, accent, bold=True, font=DISP, align=PP_ALIGN.CENTER, after=3, first=True)
    put(tf, label, 10.5, TEXT, align=PP_ALIGN.CENTER, after=0)


# ===================== SLIDE 1 — COVER =====================
s = slide()
for i, d in enumerate([5.0, 3.9, 2.8, 1.7]):
    oval_outline(s, 10.2, 3.75, d, SAFFRON if i % 2 == 0 else VIOLET)
put(tb(s, 0.9, 1.55, 8.5, 0.4), "ANANTASUTRA  ×  WONDER GIFTS", 14, SAFFRON, bold=True, first=True)
put(tb(s, 0.86, 1.95, 8.6, 1.2), "Your India Growth Partner", 46, WHITE, bold=True, font=DISP, first=True)
rect(s, 0.92, 3.0, 1.1, 0.05, SAFFRON)
put(tb(s, 0.9, 3.2, 7.8, 0.6), "Requirements discovery  ·  Our process  ·  The road ahead", 19, SAFF_LT, bold=True, first=True)
put(tb(s, 0.9, 3.95, 7.4, 1.0),
    "Bringing Dubai's premium experiential gifting to India's affluent buyers "
    "and corporates — with an accountable team on the ground.",
    14, TEXT, first=True)
put(tb(s, 0.9, 6.15, 8.5, 0.4), "14 July 2026   ·   Maria El Fassi & team, Wonder Gifts", 12.5, WHITE, bold=True, first=True)
put(tb(s, 0.9, 6.55, 8.5, 0.4), "Bhavya Duneja — Co-founder   ·   contact@anantasutra.com   ·   anantasutra.com", 12, SAFFRON, first=True)

# ===================== SLIDE 2 — TODAY'S AGENDA =====================
s = slide(); header(s, "Today", "One goal today: understand you completely.")
xs, w = cols(3)
ag = [
    ("01 · Listen", ["Your India launch priorities, target",
                     "cities and segments, current setup,",
                     "and what success looks like in 90 days."], SAFFRON),
    ("02 · Explain our process", ["How we work end-to-end — from NDA",
                                  "and proposal to execution, coordination",
                                  "and our timeline guarantee."], VIOLET),
    ("03 · Align next steps", ["Agree what happens right after this",
                               "call — so momentum never drops",
                               "while the launch window is open."], SAFFRON),
]
for x, (t, lines, acc) in zip(xs, ag):
    feature_card(s, x, 2.25, w, 3.3, t, lines, acc)
put(tb(s, 0.7, 5.85, 11.93, 0.5),
    "We don't pitch solutions before understanding problems — today is about your requirements, in detail.",
    13, SAFF_LT, italic=True, align=PP_ALIGN.CENTER, first=True)
footer(s, 2)

# ===================== SLIDE 3 — WHO WE ARE =====================
s = slide(); header(s, "Who we are", "Dedicated experts, embedded in your team.")
tf = tb(s, 0.7, 2.1, 7.5, 2.6)
put(tf, "AnantaSutra places vetted domain experts — growth, marketing, engineering, "
        "operations — inside your team. They work only for you, on your tools and "
        "your targets, and you pay just the expert's salary plus one small, "
        "transparent coordination fee. No agency markup, no lock-ins.", 15.5, TEXT, after=12, first=True)
put(tf, "Already operating across India and the UAE — your corridor is home ground for us.", 15.5, SAFF_LT, bold=True, after=0)
xs, w = cols(4)
strip = [("10+", "Brands delivered"), ("4", "Continents served"),
         ("12+", "Expert domains"), ("7-Day", "Team onboarding")]
for x, (v, l) in zip(xs, strip):
    stat_box(s, x, 5.0, w, 1.5, v, l, SAFFRON)
footer(s, 3)

# ===================== SLIDE 4 — WHAT WE UNDERSTAND SO FAR =====================
s = slide(); header(s, "Your context", "What we understand so far — you'll tell us the rest.")
xs, w = cols(3)
ctx = [
    ("The brand", ["Premium experiential gifting from Dubai —",
                   "five-star dining, adventures, staycations,",
                   "partners like Atlantis and Burj Khalifa."], SAFFRON),
    ("The move", ["A recent launch into India — a market",
                  "whose affluent buyers and corporates are",
                  "ready for exactly this product."], VIOLET),
    ("The need", ["A lead-generation engine on the ground:",
                  "HNI audiences, premium-corporate reach,",
                  "and a pipeline that runs every week."], SAFFRON),
]
for x, (t, lines, acc) in zip(xs, ctx):
    feature_card(s, x, 2.25, w, 3.3, t, lines, acc)
put(tb(s, 0.7, 5.85, 11.93, 0.5),
    "Everything on this slide is our homework — today we want to hear it in your words, with your priorities.",
    13, MUTED, italic=True, align=PP_ALIGN.CENTER, first=True)
footer(s, 4)

# ===================== SLIDE 5 — OUR PROCESS (5 STEPS) =====================
s = slide(); header(s, "Our process", "From today's call to a running engine.")
xs, w = cols(5, gap=0.24)
steps = [
    ("01", "Discovery", "Today's call — your requirements, priorities and success metrics, in detail."),
    ("02", "NDA", "We sign an NDA before anything else — your plans and data stay protected from day one."),
    ("03", "Proposal", "Every issue you raise, listed in one place — each mapped to a solution, ordered by YOUR priorities."),
    ("04", "Execution", "A definite timeline against every deliverable — with named owners and a weekly cadence."),
    ("05", "Continuity", "Weekly reviews, monthly reports, and coordinators who stay connected throughout."),
]
for x, (num, title, desc) in zip(xs, steps):
    c = card(s, x, 2.2, w, 3.6)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.22)
    put(tf, num, 28, SAFF_LT, bold=True, font=DISP, after=3, first=True)
    put(tf, title, 14.5, WHITE, bold=True, font=DISP, after=7)
    put(tf, desc, 10.5, TEXT, after=2)
put(tb(s, 0.7, 6.1, 11.93, 0.5),
    "Confidentiality first, clarity second, execution third — in that order, always.",
    13, SAFF_LT, italic=True, align=PP_ALIGN.CENTER, first=True)
footer(s, 5)

# ===================== SLIDE 6 — DEDICATED COORDINATION =====================
s = slide(); header(s, "Coordination", "Two named humans. One unbroken thread.")
xs, w = cols(2, total=8.0)
for x, (name, role) in zip(xs, [("Himanshu", "Client Coordination Lead"), ("Yash", "Client Coordination Lead")]):
    c = card(s, x, 2.3, w, 2.0, fill=CARD2)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    put(tf, name, 24, SAFF_LT, bold=True, font=DISP, align=PP_ALIGN.CENTER, after=3, first=True)
    put(tf, role, 12, TEXT, align=PP_ALIGN.CENTER, after=0)
c = card(s, 9.0, 2.3, 3.63, 2.0)
tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.24); tf.margin_right = Inches(0.2)
put(tf, "Escalation", 14, WHITE, bold=True, font=DISP, align=PP_ALIGN.CENTER, after=4, first=True)
put(tf, "Co-founder's direct line — always open to you.", 11.5, TEXT, align=PP_ALIGN.CENTER, after=0)
tf = tb(s, 0.7, 4.7, 12.0, 1.8)
put(tf, "Himanshu and Yash stay connected with your team end-to-end — WhatsApp, email or calls, "
        "in your working hours. They own the communication: requirement clarifications, progress "
        "updates, weekly reviews and monthly reports all flow through one accountable channel, "
        "so nothing ever sits in an inbox waiting.", 14.5, TEXT, first=True)
footer(s, 6)

# ===================== SLIDE 7 — TIMELINE & DELAY GUARANTEE =====================
s = slide(); header(s, "Our commitment", "A definite timeline — and we pay for our own delays.")
tf = tb(s, 0.7, 2.05, 7.3, 1.4)
put(tf, "Every proposal ships with a definite, dated timeline against each deliverable. "
        "Till date, we have never missed one.", 15.5, TEXT, after=8, first=True)
put(tf, "And we hold ourselves to it with a policy no agency offers:", 15.5, SAFF_LT, bold=True, after=0)
xs, w = cols(3)
pol = [
    ("If we ever delay…", ["Say a deliverable slips by 10 days", "from the committed timeline."], SAFFRON),
    ("…the delay is free…", ["Those next 10 days of work", "are completely free — obviously."], VIOLET),
    ("…and we credit you back.", ["PLUS 10 days of the previous month's",
                                  "fee is credited back to you.", "We pay double for our own delays."], SAFFRON),
]
for x, (t, lines, acc) in zip(xs, pol):
    feature_card(s, x, 3.6, w, 2.5, t, lines, acc, title_color=SAFF_LT, tsize=15, bsize=11.5)
put(tb(s, 0.7, 6.3, 11.93, 0.5),
    "We can afford this guarantee for one reason: we have never had to use it.",
    13.5, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
footer(s, 7)

# ===================== SLIDE 8 — PROOF =====================
s = slide(); header(s, "Proof, not promises", "Results our embedded teams have delivered.")
xs, w = cols(4)
res = [
    ("₹23L", "BotWot", "Average deal size closed across India, UAE & Nigeria — cross-border marketing, your exact corridor.", VIOLET),
    ("₹25L → ₹40L", "Awish Clinic", "Monthly revenue in 2 months — embedded web, CRM & ads experts (Delhi).", SAFFRON),
    ("+46%", "Zoom Wheels", "Revenue growth from a dedicated performance-ads expert (Delhi).", SAFFRON),
    ("+38%", "Royal Properties", "Traffic lift from embedded funnel & landing-page experts (Real Estate).", VIOLET),
]
for x, (v, n, sub, acc) in zip(xs, res):
    c = card(s, x, 2.25, w, 3.3)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.24); tf.margin_right = Inches(0.2); tf.margin_top = Inches(0.24)
    put(tf, v, 22, acc, bold=True, font=DISP, after=4, first=True)
    put(tf, n, 13, WHITE, bold=True, after=4)
    put(tf, sub, 10.5, TEXT, after=0)
put(tb(s, 0.7, 5.85, 11.93, 0.5),
    "Also delivering for brands in Qatar, the UK and Nigeria — cross-border execution is what we do.",
    13, SAFF_LT, italic=True, align=PP_ALIGN.CENTER, first=True)
footer(s, 8)

# ===================== SLIDE 9 — THE TWO TRACKS (SHAPED TODAY) =====================
s = slide(); header(s, "Where this can go", "Two tracks — final shape comes from today's discussion.")
xs, w = cols(2)
tr = [
    ("Track 1 — Embedded Growth Team", [
        "A dedicated lead-gen / growth expert working only",
        "for Wonder Gifts — performance campaigns, corporate",
        "outreach, partnerships and CRM follow-up.",
        "",
        "Live in 7 days · you pay just the salary · 48-hour",
        "replacement guarantee · scale or pause anytime."], SAFFRON),
    ("Track 2 — HNI & Corporate Reach", [
        "Curated, segmented HNI and premium-corporate",
        "decision-maker audiences — by city, industry and",
        "spend profile — with targeted, compliant outreach.",
        "",
        "Built for an AED 300–3,000+ experiential product:",
        "reach exactly the buyers it was made for."], VIOLET),
]
for x, (t, lines, acc) in zip(xs, tr):
    feature_card(s, x, 2.25, w, 3.7, t, lines, acc, title_color=SAFF_LT, tsize=16, bsize=12)
footer(s, 9)

# ===================== SLIDE 10 — NEXT STEPS / CLOSE =====================
s = slide()
for i, d in enumerate([5.4, 4.1, 2.8]):
    oval_outline(s, 6.66, 3.1, d, SAFFRON if i % 2 == 0 else VIOLET)
put(tb(s, 1, 1.5, 11.3, 0.4), "THE ROAD FROM HERE", 14, SAFFRON, bold=True, align=PP_ALIGN.CENTER, first=True)
put(tb(s, 1, 1.95, 11.3, 1.0), "Let's build your India engine.", 40, WHITE, bold=True, font=DISP, align=PP_ALIGN.CENTER, first=True)
tf = tb(s, 2.4, 3.3, 8.5, 2.0)
put(tf, "Today — your requirements, understood in full", 15, TEXT, align=PP_ALIGN.CENTER, after=6, first=True)
put(tf, "Within 24 hours — NDA shared and signed", 15, TEXT, align=PP_ALIGN.CENTER, after=6)
put(tf, "Within the week — full proposal: every issue listed, resolved, with a definite timeline", 15, TEXT, align=PP_ALIGN.CENTER, after=6)
put(tf, "Then — Himanshu & Yash connected, execution begins", 15, SAFF_LT, bold=True, align=PP_ALIGN.CENTER, after=0)
cta = card(s, 4.92, 5.55, 3.5, 0.7, fill=SAFFRON, border=SAFFRON)
ctf = cta.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
put(ctf, "contact@anantasutra.com", 15, BG, bold=True, align=PP_ALIGN.CENTER, first=True)
put(tb(s, 1, 6.5, 11.3, 0.4), "Bhavya Duneja — Co-founder   ·   anantasutra.com   ·   Delhi, India", 12, MUTED, align=PP_ALIGN.CENTER, first=True)

out = "AnantaSutra-WonderGifts-Deck.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
